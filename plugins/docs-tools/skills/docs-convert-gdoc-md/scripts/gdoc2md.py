"""
Export Google Docs to Markdown, Slides to Markdown (via PPTX),
or Sheets to CSV.  Optionally include Google Docs comments as
Markdown footnotes.

Requires gcloud CLI and python-pptx (for Slides export).

python3 ${CLAUDE_SKILL_DIR}/scripts/gdoc2md.py [--comments] [--include-resolved] <url> [output]
"""

import argparse
import json
import re
import subprocess
import sys
import time
from io import BytesIO
from pathlib import Path
from urllib.error import HTTPError
from urllib.parse import quote
from urllib.request import Request, urlopen

# tolerates trailing segments like /edit, /view, ?usp=sharing
VALID_URL_RE = re.compile(
    r"^https://docs\.google\.com/"
    r"(?P<mode>document|presentation|spreadsheets)/d/(?P<id>[a-zA-Z0-9_-]+)"
)

MODE_MAP = {
    "document": "doc",
    "presentation": "slides",
    "spreadsheets": "sheets",
}

EXTENSIONS = {"doc": ".md", "slides": ".md", "sheets": ".csv"}


# ---------------------------------------------------------------------------
# Argument parsing & validation
# ---------------------------------------------------------------------------


def parse_and_validate_args():
    """Parse CLI arguments and return validated args tuple."""
    parser = argparse.ArgumentParser(
        description="Export Google Docs/Slides/Sheets to Markdown or CSV.",
    )
    parser.add_argument("url", help="Google Docs, Slides, or Sheets URL")
    parser.add_argument("output", nargs="?", default=None, help="Output file path")
    parser.add_argument(
        "--comments",
        action="store_true",
        help="Include Google Docs comments as Markdown footnotes (Docs only)",
    )
    parser.add_argument(
        "--include-resolved",
        action="store_true",
        help="Include resolved comment threads (requires --comments)",
    )
    parser.add_argument(
        "--manifest",
        action="store_true",
        help="Write a companion manifest file with section map (Docs only)",
    )
    parser.add_argument(
        "--split-sections",
        action="store_true",
        help=(
            "Split output into per-section files under 40 KB each (requires --manifest, Docs only)"
        ),
    )
    args = parser.parse_args()

    if args.include_resolved and not args.comments:
        parser.error("--include-resolved requires --comments")

    if args.split_sections and not args.manifest:
        parser.error("--split-sections requires --manifest")

    match = VALID_URL_RE.match(args.url)
    if not match:
        parser.error(
            "URL must be a Google Docs, Slides, or Sheets URL (https://docs.google.com/...)"
        )

    mode = MODE_MAP[match.group("mode")]
    file_id = match.group("id")
    output = args.output or f"{file_id}{EXTENSIONS[mode]}"

    if args.comments and mode != "doc":
        print(
            "Warning: --comments is only supported for Google Docs, ignoring.",
            file=sys.stderr,
        )

    return (
        file_id,
        output,
        mode,
        args.comments,
        args.include_resolved,
        args.manifest,
        args.split_sections,
    )


# ---------------------------------------------------------------------------
# Dependency check
# ---------------------------------------------------------------------------


def check_dependencies():
    """Verify that the gcloud CLI is installed, exiting with guidance if not."""
    result = subprocess.run(["gcloud", "version"], capture_output=True)  # noqa: S607
    if result.returncode != 0:
        print("Error: gcloud CLI is not installed.", file=sys.stderr)
        print(
            "  Install: https://cloud.google.com/sdk/docs/install",
            file=sys.stderr,
        )
        sys.exit(1)


# ---------------------------------------------------------------------------
# Auth — single source of truth for obtaining a token
# ---------------------------------------------------------------------------


def get_token() -> str:
    """
    Return a valid access token, prompting the user to log in if needed.
    Raises SystemExit on unrecoverable failure.
    """
    result = subprocess.run(
        ["gcloud", "auth", "print-access-token"],  # noqa: S607
        capture_output=True,
        text=True,
    )
    if result.returncode == 0 and result.stdout.strip():
        return result.stdout.strip()

    print("No active credentials found. Authenticating with Google...")
    login = subprocess.run(
        ["gcloud", "auth", "login", "--enable-gdrive-access"],  # noqa: S607
    )
    if login.returncode != 0:
        print(
            "Error: Authentication was cancelled or failed.",
            file=sys.stderr,
        )
        sys.exit(1)

    # Re-fetch after successful login
    result = subprocess.run(
        ["gcloud", "auth", "print-access-token"],  # noqa: S607
        capture_output=True,
        text=True,
    )
    if result.returncode != 0 or not result.stdout.strip():
        print(
            "Error: Could not obtain access token after login.",
            file=sys.stderr,
        )
        sys.exit(1)

    return result.stdout.strip()


# ---------------------------------------------------------------------------
# Download
# ---------------------------------------------------------------------------


def download(url: str, token: str, retries: int = 3) -> bytes:
    """GET *url* with Bearer auth and exponential back-off on 429 responses."""
    req = Request(url, headers={"Authorization": f"Bearer {token}"})  # noqa: S310
    for attempt in range(retries + 1):
        try:
            with urlopen(req) as resp:  # noqa: S310
                return resp.read()
        except HTTPError as e:
            if e.code == 429 and attempt < retries:
                wait = 2**attempt
                print(
                    f"Rate limited (429), retrying in {wait}s...",
                    file=sys.stderr,
                )
                time.sleep(wait)
                continue
            messages = {
                401: ("Authentication failed (401). Try: gcloud auth login --enable-gdrive-access"),
                403: ("Access denied (403). Check you have permission to access this file."),
                404: "Not found (404). Check the URL is correct.",
            }
            print(
                f"Error: {messages.get(e.code, f'HTTP {e.code}')}",
                file=sys.stderr,
            )
            sys.exit(1)


# ---------------------------------------------------------------------------
# PPTX → Markdown conversion
# ---------------------------------------------------------------------------


def pptx_to_markdown(data: bytes) -> str:
    """
    Convert PPTX bytes to structured Markdown.

    Note: images, charts, and other non-text/non-table shapes are not exported.
    Only text frames and tables are extracted.
    """
    try:
        from pptx import Presentation
    except ImportError:
        print(
            "Error: python-pptx is required for Slides export.",
            file=sys.stderr,
        )
        print("  Install: pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(BytesIO(data))
    lines = []

    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"## Slide {i}")
        lines.append("")

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.replace("\x0b", "\n").strip()
                    if not text:
                        continue
                    if paragraph.level > 0:
                        indent = "  " * (paragraph.level - 1)
                        for subline in text.split("\n"):
                            lines.append(f"{indent}- {subline}")
                    else:
                        lines.append(text)
                lines.append("")

            elif shape.has_table:
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    cells = [cell.text.strip().replace("|", "\\|") for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        sep = "| " + " | ".join(["---"] * len(cells)) + " |"
                        lines.append(sep)
                lines.append("")

        # Images, charts, and other shape types are not exported
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                lines.append("> **Notes:** " + notes_text)
                lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Sheets metadata
# ---------------------------------------------------------------------------


def get_sheet_metadata(file_id: str, token: str):
    """Return list of (gid, title) for each sheet."""
    api_url = f"https://sheets.googleapis.com/v4/spreadsheets/{file_id}?fields=sheets.properties"
    data = download(api_url, token)
    info = json.loads(data)
    return [
        (
            s["properties"]["sheetId"],
            s["properties"]["title"],
        )
        for s in info["sheets"]
    ]


def _sanitize_filename(name: str) -> str:
    """Replace filesystem-unsafe characters with underscores."""
    return re.sub(r'[\\/*?:"<>|]', "_", name)


# ---------------------------------------------------------------------------
# Google Docs comments → Markdown footnotes
# ---------------------------------------------------------------------------


def fetch_comments(
    file_id: str,
    token: str,
    include_resolved: bool = False,
) -> list[dict]:
    """Fetch comment threads from the Drive v3 API.

    Returns a list of dicts with keys: author, content, quoted_text,
    resolved, and replies (list of {author, content}).
    """
    fields = (
        "nextPageToken,"
        "comments(id,content,resolved,author/displayName,"
        "quotedFileContent/value,replies(content,author/displayName))"
    )
    comments = []
    page_token = None
    while True:
        api_url = (
            f"https://www.googleapis.com/drive/v3/files/{file_id}/comments"
            f"?fields={quote(fields, safe='()/,')}&includeDeleted=false"
            f"&pageSize=100"
        )
        if page_token:
            api_url += f"&pageToken={quote(page_token)}"

        data = json.loads(download(api_url, token))
        for c in data.get("comments", []):
            resolved = c.get("resolved", False)
            if resolved and not include_resolved:
                continue
            quoted = (c.get("quotedFileContent") or {}).get("value", "")
            replies = [
                {
                    "author": r.get("author", {}).get("displayName", "Unknown"),
                    "content": r.get("content", ""),
                }
                for r in c.get("replies", [])
            ]
            comments.append(
                {
                    "author": c.get("author", {}).get("displayName", "Unknown"),
                    "content": c.get("content", ""),
                    "quoted_text": quoted,
                    "resolved": resolved,
                    "replies": replies,
                }
            )

        page_token = data.get("nextPageToken")
        if not page_token:
            break

    return comments


def _normalize(text: str) -> str:
    """Collapse whitespace for fuzzy anchor matching."""
    return re.sub(r"\s+", " ", text).strip()


def insert_comment_footnotes(
    markdown: str,
    comments: list[dict],
) -> str:
    """Insert footnote references into the Markdown body and append
    footnote definitions at the end of the file.

    Matching strategy: for each comment with a quoted anchor, find the
    first occurrence of that anchor text in the Markdown (normalized
    whitespace) and insert a footnote reference after it.  Comments
    without an anchor are appended as unanchored footnotes at the end.
    """
    if not comments:
        return markdown

    footnotes: list[str] = []
    fn_index = 1

    # Pass 1: resolve all anchor positions against the *unmodified* markdown
    # so earlier matches cannot invalidate later ones.
    norm_md = _normalize(markdown)
    used_offsets: set[int] = set()
    insertions: list[tuple[int, str, str]] = []

    for comment in comments:
        anchor = comment["quoted_text"]
        label = f"[^{fn_index}]"

        body_parts = []
        status = " (resolved)" if comment["resolved"] else ""
        body_parts.append(f"**{comment['author']}{status}:** {_normalize(comment['content'])}")
        for reply in comment["replies"]:
            body_parts.append(f"    **{reply['author']}:** {_normalize(reply['content'])}")
        footnote_def = f"{label}: " + " \\\n".join(body_parts)

        norm_anchor = _normalize(anchor) if anchor else ""
        if norm_anchor:
            search_from = 0
            pos = -1
            while True:
                candidate = norm_md.find(norm_anchor, search_from)
                if candidate == -1:
                    break
                orig_end = _find_original_end(markdown, norm_md, candidate, len(norm_anchor))
                if orig_end not in used_offsets:
                    pos = candidate
                    break
                search_from = candidate + 1
            if pos != -1:
                end_of_anchor = _find_original_end(markdown, norm_md, pos, len(norm_anchor))
                end_of_anchor = _snap_to_word_boundary(markdown, end_of_anchor)
                used_offsets.add(end_of_anchor)
                insertions.append((end_of_anchor, label, footnote_def))
                fn_index += 1
                continue

        footnotes.append(footnote_def)
        fn_index += 1

    # Pass 2: apply insertions from end to start so offsets stay valid.
    insertions.sort(key=lambda t: t[0], reverse=True)
    for offset, label, footnote_def in insertions:
        markdown = markdown[:offset] + label + markdown[offset:]
        footnotes.append(footnote_def)

    # Re-sort footnotes by their numeric index for consistent output.
    footnotes.sort(key=lambda f: int(f.split("]")[0].lstrip("[^")))

    if footnotes:
        markdown = markdown.rstrip() + "\n\n---\n\n"
        markdown += "\n".join(footnotes) + "\n"

    return markdown


def _find_original_end(
    original: str,
    normalized: str,
    norm_pos: int,
    norm_len: int,
) -> int:
    """Map a position in the normalized string back to the original.

    Walk through the original string, tracking how many non-collapsed
    characters have been consumed, to find where the anchor ends in
    the original text.
    """
    consumed = 0
    i = 0
    in_space = False

    while i < len(original) and consumed < norm_pos:
        if original[i].isspace():
            if not in_space:
                consumed += 1
                in_space = True
        else:
            consumed += 1
            in_space = False
        i += 1

    chars_left = norm_len
    while i < len(original) and chars_left > 0:
        if original[i].isspace():
            if not in_space:
                chars_left -= 1
                in_space = True
        else:
            chars_left -= 1
            in_space = False
        i += 1

    return i


def _snap_to_word_boundary(text: str, pos: int) -> int:
    """Advance *pos* past any remaining word characters so the footnote
    reference never splits a word.  Stops at whitespace, punctuation
    that commonly follows words, or end-of-string.
    """
    while pos < len(text) and text[pos].isalnum():
        pos += 1
    return pos


# ---------------------------------------------------------------------------
# Manifest generation
# ---------------------------------------------------------------------------

HEADING_RE = re.compile(r"^#{1,3}\s")
SUB_HEADING_RE = re.compile(r"^#{4,6}\s")


def _extract_brief(lines: list[str], max_len: int = 120) -> str:
    """Return the first non-heading, non-blank line as a content brief, truncated to max_len."""
    for line in lines:
        stripped = line.strip()
        if not stripped:
            continue
        if HEADING_RE.match(stripped) or SUB_HEADING_RE.match(stripped):
            continue
        if stripped.startswith(("---", "===", "```", "<!--", "|", "![")):
            continue
        if re.match(r"^[-*]\s|^\d+\.\s", stripped):
            continue
        if len(stripped) > max_len:
            return stripped[: max_len - 1] + "…"
        return stripped
    return ""


def generate_manifest(
    content: str, output_path: str, section_files: list[dict] | None = None
) -> dict:
    lines = content.splitlines()
    total_lines = len(lines)

    title = Path(output_path).stem
    for line in lines:
        if HEADING_RE.match(line):
            title = line.lstrip("#").strip()
            break

    sections: list[tuple[str, int, int]] = []
    for i, line in enumerate(lines):
        if HEADING_RE.match(line):
            sections.append((line.strip(), i, -1))

    section_entries: list[str] = []
    for idx, (heading, start, _) in enumerate(sections):
        end = sections[idx + 1][1] - 1 if idx + 1 < len(sections) else total_lines - 1
        char_count = sum(len(lines[j]) for j in range(start, end + 1))
        section_entries.append(
            f"- Line {start + 1}-{end + 1}: {heading.lstrip('#').strip()} ({char_count:,} chars)"
        )

    total_chars = len(content)
    total_estimated_tokens = total_chars // 3

    manifest_lines = [
        f"# Document Manifest: {title}",
        f"Total characters: {total_chars}",
        f"Total estimated tokens: {total_estimated_tokens}",
        "",
        "## Sections",
    ]
    manifest_lines.extend(section_entries)

    if section_files:
        manifest_lines.append("")
        manifest_lines.append("## Section Files")
        for sf in section_files:
            brief = sf.get("brief", "")
            brief_suffix = f" — {brief}" if brief else ""
            manifest_lines.append(
                f"- {Path(sf['file']).name}: {sf['heading']} ({sf['chars']:,} chars){brief_suffix}"
            )

    manifest_file = f"{output_path}.manifest.md"
    Path(manifest_file).write_text("\n".join(manifest_lines) + "\n", encoding="utf-8")

    result = {
        "manifest_file": manifest_file,
        "total_chars": total_chars,
        "total_estimated_tokens": total_estimated_tokens,
        "section_count": len(sections),
    }
    if section_files:
        result["section_files"] = section_files
    return result


def _split_at_blank_lines(lines: list[str], max_bytes: int) -> list[list[str]]:
    """Split a list of lines into chunks at blank-line boundaries, each under max_bytes."""
    chunks: list[list[str]] = []
    current: list[str] = []
    current_bytes = 0
    for line in lines:
        line_bytes = len(line.encode("utf-8")) + 1
        if current_bytes + line_bytes > max_bytes and current:
            if line.strip() == "":
                chunks.append(current)
                current = []
                current_bytes = 0
            elif current_bytes >= max_bytes:
                chunks.append(current)
                current = []
                current_bytes = 0
        current.append(line)
        current_bytes += line_bytes
    if current:
        chunks.append(current)
    return chunks


def _split_oversized_section(
    section_lines: list[str],
    parent: Path,
    stem: str,
    section_idx: int,
    max_bytes: int,
) -> list[dict]:
    """Split a section that exceeds max_bytes into sub-chunks."""
    sub_boundaries = [0]
    for i, line in enumerate(section_lines):
        if i > 0 and SUB_HEADING_RE.match(line):
            sub_boundaries.append(i)

    if len(sub_boundaries) > 1:
        sub_boundaries.append(len(section_lines))
        raw_chunks = []
        for j in range(len(sub_boundaries) - 1):
            start, end = sub_boundaries[j], sub_boundaries[j + 1]
            raw_chunks.append(section_lines[start:end])
        chunks = []
        for rc in raw_chunks:
            rc_bytes = sum(len(line.encode("utf-8")) + 1 for line in rc)
            if rc_bytes <= max_bytes:
                chunks.append(rc)
            else:
                sub = _split_at_blank_lines(rc, max_bytes)
                chunks.extend(sub)
    else:
        chunks = _split_at_blank_lines(section_lines, max_bytes)

    results = []
    for j, chunk in enumerate(chunks):
        suffix = chr(97 + j) if j < 26 else str(j)
        filename = f"{stem}-section-{section_idx:02d}{suffix}.md"
        filepath = parent / filename
        text = "\n".join(chunk)
        filepath.write_text(text + "\n", encoding="utf-8")
        heading = (
            chunk[0].lstrip("#").strip()
            if HEADING_RE.match(chunk[0]) or SUB_HEADING_RE.match(chunk[0])
            else f"(continued part {suffix})"
        )
        results.append(
            {
                "file": str(filepath),
                "heading": heading,
                "chars": len(text),
                "brief": _extract_brief(chunk),
            }
        )
    return results


def split_into_section_files(
    content: str, output_path: str, max_section_bytes: int = 40000
) -> list[dict]:
    """Split markdown content into per-section files, each under max_section_bytes."""
    lines = content.splitlines()
    total_lines = len(lines)
    parent = Path(output_path).parent
    stem = Path(output_path).stem

    boundaries = []
    for i, line in enumerate(lines):
        if HEADING_RE.match(line):
            boundaries.append(i)

    if not boundaries:
        filename = f"{stem}-section-01.md"
        filepath = parent / filename
        filepath.write_text(content + "\n", encoding="utf-8")
        return [
            {
                "file": str(filepath),
                "heading": stem,
                "chars": len(content),
                "brief": _extract_brief(lines),
            }
        ]

    if boundaries[0] > 0:
        boundaries.insert(0, 0)

    results = []
    section_idx = 1
    for b_idx in range(len(boundaries)):
        start = boundaries[b_idx]
        end = boundaries[b_idx + 1] if b_idx + 1 < len(boundaries) else total_lines
        section_lines = lines[start:end]
        section_text = "\n".join(section_lines)
        section_bytes = len(section_text.encode("utf-8"))

        if section_bytes <= max_section_bytes:
            filename = f"{stem}-section-{section_idx:02d}.md"
            filepath = parent / filename
            filepath.write_text(section_text + "\n", encoding="utf-8")
            heading = (
                section_lines[0].lstrip("#").strip() if HEADING_RE.match(section_lines[0]) else stem
            )
            results.append(
                {
                    "file": str(filepath),
                    "heading": heading,
                    "chars": len(section_text),
                    "brief": _extract_brief(section_lines),
                }
            )
            section_idx += 1
        else:
            sub_results = _split_oversized_section(
                section_lines, parent, stem, section_idx, max_section_bytes
            )
            results.extend(sub_results)
            section_idx += 1

    return results


# ---------------------------------------------------------------------------
# Fetch & write
# ---------------------------------------------------------------------------


def fetch(
    file_id: str,
    output: str,
    mode: str,
    include_comments: bool = False,
    include_resolved: bool = False,
):
    """Download and convert a Google Docs/Slides/Sheets file, writing the result to *output*."""
    token = get_token()
    base = "https://docs.google.com"

    if mode == "sheets":
        _fetch_sheets(file_id, output, token, base)
        return

    export_urls = {
        "slides": (f"{base}/presentation/d/{file_id}/export?format=pptx"),
        "doc": (f"{base}/document/d/{file_id}/export?format=md"),
    }

    data = download(export_urls[mode], token)
    output_path = Path(output)

    if output_path.exists():
        print(
            f"Warning: overwriting existing file '{output}'",
            file=sys.stderr,
        )

    if mode == "slides":
        output_path.write_text(pptx_to_markdown(data), encoding="utf-8")
    elif mode == "doc":
        md_text = data.decode("utf-8")
        if include_comments:
            comments = fetch_comments(
                file_id,
                token,
                include_resolved,
            )
            if comments:
                md_text = insert_comment_footnotes(md_text, comments)
                print(
                    f"Inserted {len(comments)} comment(s) as footnotes.",
                    file=sys.stderr,
                )
            else:
                print("No comments found.", file=sys.stderr)
        output_path.write_text(md_text, encoding="utf-8")
    else:
        output_path.write_bytes(data)

    print(f"Saved to {output}")


def _fetch_sheets(file_id: str, output: str, token: str, base: str):
    """Export every sheet in a spreadsheet as a separate CSV."""
    try:
        sheets = get_sheet_metadata(file_id, token)
    except SystemExit:
        # Sheets API not enabled — fall back to default first sheet
        print(
            "Warning: Could not fetch sheet metadata "
            "(Sheets API may not be enabled). "
            "Exporting first sheet only.",
            file=sys.stderr,
        )
        sheets = [(0, "Sheet1")]
    out_path = Path(output)
    stem = out_path.stem
    parent = out_path.parent

    if len(sheets) == 1:
        gid, title = sheets[0]
        url = f"{base}/spreadsheets/d/{file_id}/export?format=csv&gid={gid}"
        data = download(url, token)
        if out_path.exists():
            print(
                f"Warning: overwriting existing file '{out_path}'",
                file=sys.stderr,
            )
        out_path.write_bytes(data)
        print(f"Saved to {out_path}")
        return

    for gid, title in sheets:
        safe_title = _sanitize_filename(title)
        csv_path = parent / f"{stem}_{safe_title}.csv"
        url = f"{base}/spreadsheets/d/{file_id}/export?format=csv&gid={gid}"
        data = download(url, token)
        if csv_path.exists():
            print(
                f"Warning: overwriting existing file '{csv_path}'",
                file=sys.stderr,
            )
        csv_path.write_bytes(data)
        print(f"Saved to {csv_path}")


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------


def main():
    """CLI entry point: parse arguments, check dependencies, and run the export."""
    file_id, output, mode, comments, include_resolved, manifest, split_sections = (
        parse_and_validate_args()
    )
    check_dependencies()
    fetch(file_id, output, mode, comments, include_resolved)

    if manifest and mode == "doc":
        content = Path(output).read_text(encoding="utf-8")
        section_files = None
        if split_sections:
            section_files = split_into_section_files(content, output)
        result = generate_manifest(content, output, section_files)
        result["output_file"] = output
        print(json.dumps(result))


if __name__ == "__main__":
    main()
